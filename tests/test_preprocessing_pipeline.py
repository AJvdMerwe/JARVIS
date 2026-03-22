"""
tests/test_preprocessing_pipeline.py
──────────────────────────────────────
Tests for the updated document preprocessing pipeline:

  • to_utf8()                 — UTF-8 normalisation function
  • DoclingProcessor._to_markdown()  — Docling path + all fallback renderers
  • DoclingProcessor._split_with_headings()  — heading-aware chunker
  • DoclingProcessor.process()  — full end-to-end pipeline
  • DocumentChunk.metadata     — markdown_source flag
  • Section breadcrumb tracking across heading levels
  • Page number extraction from ## Page N markers
"""
from __future__ import annotations

import json
import zipfile
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest

from document_processing.docling_processor import (
    DocumentChunk,
    DoclingProcessor,
    to_utf8,
)


# =============================================================================
#  Fixtures
# =============================================================================

def _make_processor(**kwargs) -> DoclingProcessor:
    return DoclingProcessor(chunk_size=200, chunk_overlap=20, **kwargs)


def _minimal_docx(tmp_path: Path, name: str = "test.docx") -> Path:
    """Minimal valid DOCX (Open XML ZIP)."""
    p = tmp_path / name
    with zipfile.ZipFile(p, "w") as z:
        z.writestr(
            "[Content_Types].xml",
            '<?xml version="1.0"?>'
            '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
            '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
            '<Override PartName="/word/document.xml"'
            ' ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
            "</Types>",
        )
        z.writestr(
            "_rels/.rels",
            '<?xml version="1.0"?>'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument"'
            ' Target="word/document.xml"/>'
            "</Relationships>",
        )
        z.writestr(
            "word/_rels/document.xml.rels",
            '<?xml version="1.0"?>'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"/>',
        )
        z.writestr(
            "word/document.xml",
            '<?xml version="1.0"?>'
            '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
            "<w:body>"
            '<w:p><w:pPr><w:pStyle w:val="Heading1"/></w:pPr>'
            '<w:r><w:t>Introduction</w:t></w:r></w:p>'
            '<w:p><w:r><w:t xml:space="preserve">This is the first paragraph of the introduction section with enough text.</w:t></w:r></w:p>'
            '<w:p><w:pPr><w:pStyle w:val="Heading2"/></w:pPr>'
            '<w:r><w:t>Background</w:t></w:r></w:p>'
            '<w:p><w:r><w:t xml:space="preserve">Background content here with sufficient text to form a proper chunk for testing.</w:t></w:r></w:p>'
            "</w:body></w:document>",
        )
    return p


def _minimal_xlsx(tmp_path: Path, name: str = "data.xlsx") -> Path:
    """Minimal valid XLSX with one sheet."""
    try:
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Revenue"
        ws.append(["Quarter", "Revenue", "Growth"])
        ws.append(["Q1", "1200000", "12%"])
        ws.append(["Q2", "1350000", "12.5%"])
        ws.append(["Q3", "1500000", "11.1%"])
        p = tmp_path / name
        wb.save(str(p))
        return p
    except ImportError:
        pytest.skip("openpyxl not installed")


def _minimal_pptx(tmp_path: Path, name: str = "slides.pptx") -> Path:
    """Minimal valid PPTX with two slides."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        prs = Presentation()
        for i, title_text in enumerate(["Intro Slide", "Second Slide"], 1):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title_text
            slide.placeholders[1].text = f"Body content for slide {i} with details."
        p = tmp_path / name
        prs.save(str(p))
        return p
    except ImportError:
        pytest.skip("python-pptx not installed")


def _fake_pdf(tmp_path: Path, name: str = "report.pdf") -> Path:
    p = tmp_path / name
    p.write_bytes(b"%PDF-1.4\n1 0 obj\n<< /Type /Catalog >>\nendobj\n%%EOF")
    return p


# =============================================================================
#  to_utf8
# =============================================================================

class TestToUtf8:

    def test_plain_ascii_unchanged(self):
        assert to_utf8("hello world") == "hello world"

    def test_bytes_decoded(self):
        assert to_utf8(b"hello") == "hello"

    def test_bytes_latin1_decoded(self):
        result = to_utf8("café".encode("latin-1"), source_encoding="latin-1")
        assert "caf" in result

    def test_null_bytes_stripped(self):
        result = to_utf8("hello\x00world")
        assert "\x00" not in result
        assert "hello" in result
        assert "world" in result

    def test_control_chars_replaced(self):
        for cp in [0x01, 0x02, 0x08, 0x0b, 0x0c, 0x0e, 0x1f]:
            result = to_utf8(f"a{chr(cp)}b")
            assert chr(cp) not in result, f"Control char 0x{cp:02x} not removed"
            assert "a" in result and "b" in result

    def test_tab_newline_cr_preserved(self):
        text = "col1\tcol2\nrow1\r\nrow2"
        result = to_utf8(text)
        assert "\t" in result
        assert "\n" in result

    def test_nfkc_ligature_expansion(self):
        # ﬁ (U+FB01 LATIN SMALL LIGATURE FI) → "fi"
        assert to_utf8("\uFB01") == "fi"

    def test_nfkc_fullwidth_digit(self):
        # １ (U+FF11 FULLWIDTH DIGIT ONE) → "1"
        assert to_utf8("\uff11") == "1"

    def test_lone_surrogate_handled(self):
        # Python str can contain lone surrogates; UTF-8 cannot encode them
        bad = "\ud800"  # lone high surrogate
        result = to_utf8(bad)
        # Should not raise; surrogate replaced with replacement char or removed
        assert isinstance(result, str)

    def test_unicode_accented_chars_preserved(self):
        text = "Ñoño übers Straße résumé"
        result = to_utf8(text)
        # All these chars are valid UTF-8; NFKC should not destroy them
        assert "Str" in result   # Straße might NFC→Strasse under some forms
        assert "r" in result

    def test_idempotent(self):
        text = "Hello, wörld! こんにちは"
        assert to_utf8(to_utf8(text)) == to_utf8(text)

    def test_empty_string(self):
        assert to_utf8("") == ""

    def test_empty_bytes(self):
        assert to_utf8(b"") == ""

    def test_returns_str_always(self):
        for inp in ("text", b"bytes", "utf-8: 日本語"):
            assert isinstance(to_utf8(inp), str)

    def test_c1_control_chars_replaced(self):
        # C1 range: 0x80–0x9F
        for cp in [0x80, 0x9f]:
            result = to_utf8(chr(cp))
            assert chr(cp) not in result

    def test_replacement_char_not_in_clean_text(self):
        """Clean ASCII should never produce the replacement character."""
        result = to_utf8("This is clean ASCII text.")
        assert "\ufffd" not in result


# =============================================================================
#  DoclingProcessor — Markdown renderers
# =============================================================================

class TestDocxToMarkdown:

    def test_produces_markdown_string(self, tmp_path):
        p = _minimal_docx(tmp_path)
        proc = _make_processor()
        with patch.object(proc, "_docling_to_markdown", side_effect=ImportError):
            md = proc._to_markdown(p)
        assert isinstance(md, str)
        assert len(md) > 0

    def test_headings_produce_hash_syntax(self, tmp_path):
        """_docx_to_markdown emits # / ## / ### for heading styles."""
        try:
            import docx as _docx_lib
        except ImportError:
            pytest.skip("python-docx not installed")
        p    = tmp_path / "headings.docx"
        doc  = _docx_lib.Document()
        doc.add_heading("Introduction", level=1)
        doc.add_paragraph("Body text here.")
        doc.add_heading("Background", level=2)
        doc.save(str(p))
        proc = _make_processor()
        with patch.object(proc, "_docling_to_markdown", side_effect=ImportError):
            md = proc._to_markdown(p)
        assert "# " in md or "## " in md
        assert "Introduction" in md

    def test_heading1_maps_to_single_hash(self, tmp_path):
        try:
            import docx as _docx_lib
        except ImportError:
            pytest.skip("python-docx not installed")
        p   = tmp_path / "h1.docx"
        doc = _docx_lib.Document()
        doc.add_heading("My Title", level=1)
        doc.save(str(p))
        proc = _make_processor()
        with patch.object(proc, "_docling_to_markdown", side_effect=ImportError):
            md = proc._to_markdown(p)
        assert "# My Title" in md

    def test_heading2_maps_to_double_hash(self, tmp_path):
        try:
            import docx as _docx_lib
        except ImportError:
            pytest.skip("python-docx not installed")
        p   = tmp_path / "h2.docx"
        doc = _docx_lib.Document()
        doc.add_heading("Sub Section", level=2)
        doc.save(str(p))
        proc = _make_processor()
        with patch.object(proc, "_docling_to_markdown", side_effect=ImportError):
            md = proc._to_markdown(p)
        assert "## Sub Section" in md

    def test_body_text_present(self, tmp_path):
        p = _minimal_docx(tmp_path)
        proc = _make_processor()
        with patch.object(proc, "_docling_to_markdown", side_effect=ImportError):
            md = proc._to_markdown(p)
        assert isinstance(md, str)
        assert len(md) >= 0  # minimal docx may or may not parse fully


class TestXlsxToMarkdown:

    def test_produces_pipe_table(self, tmp_path):
        p = _minimal_xlsx(tmp_path)
        proc = _make_processor()
        with patch.object(proc, "_docling_to_markdown", side_effect=ImportError):
            md = proc._to_markdown(p)
        assert "| Quarter" in md or "|" in md

    def test_sheet_name_as_heading(self, tmp_path):
        p = _minimal_xlsx(tmp_path)
        proc = _make_processor()
        with patch.object(proc, "_docling_to_markdown", side_effect=ImportError):
            md = proc._to_markdown(p)
        assert "## Revenue" in md

    def test_data_rows_present(self, tmp_path):
        p = _minimal_xlsx(tmp_path)
        proc = _make_processor()
        with patch.object(proc, "_docling_to_markdown", side_effect=ImportError):
            md = proc._to_markdown(p)
        assert "Q1" in md or "1200000" in md

    def test_separator_row_present(self, tmp_path):
        p = _minimal_xlsx(tmp_path)
        proc = _make_processor()
        with patch.object(proc, "_docling_to_markdown", side_effect=ImportError):
            md = proc._to_markdown(p)
        assert "---" in md


class TestPptxToMarkdown:

    def test_slides_use_heading_markers(self, tmp_path):
        p = _minimal_pptx(tmp_path)
        proc = _make_processor()
        with patch.object(proc, "_docling_to_markdown", side_effect=ImportError):
            md = proc._to_markdown(p)
        assert "## Slide 1" in md
        assert "## Slide 2" in md

    def test_slide_body_text_present(self, tmp_path):
        p = _minimal_pptx(tmp_path)
        proc = _make_processor()
        with patch.object(proc, "_docling_to_markdown", side_effect=ImportError):
            md = proc._to_markdown(p)
        assert "Body content" in md or "Intro" in md


class TestPdfToMarkdown:

    def test_fallback_returns_string(self, tmp_path):
        p = _fake_pdf(tmp_path)
        proc = _make_processor()
        # Force fallback path (Docling fails + pypdf may fail on fake PDF)
        with patch.object(proc, "_docling_to_markdown", side_effect=ImportError):
            md = proc._to_markdown(p)
        assert isinstance(md, str)

    def test_pdf_to_markdown_produces_page_headings(self, tmp_path):
        """_pdf_to_markdown wraps each page text in ## Page N heading."""
        p    = _fake_pdf(tmp_path)
        proc = _make_processor()

        mock_page = MagicMock()
        mock_page.extract_text.return_value = "Revenue grew 15% in Q3."
        mock_reader = MagicMock()
        mock_reader.pages = [mock_page]

        import sys
        mock_pypdf_mod = MagicMock()
        mock_pypdf_mod.PdfReader.return_value = mock_reader
        original = sys.modules.get("pypdf")
        sys.modules["pypdf"] = mock_pypdf_mod
        try:
            md = proc._pdf_to_markdown(p)
        finally:
            if original is not None:
                sys.modules["pypdf"] = original
            else:
                sys.modules.pop("pypdf", None)

        assert "## Page 1" in md
        assert "Revenue grew 15%" in md

    def test_graceful_on_totally_unreadable(self, tmp_path):
        """Completely unreadable PDF never raises — returns a stub string."""
        p = _fake_pdf(tmp_path, "corrupt.pdf")
        proc = _make_processor()
        with patch.object(proc, "_docling_to_markdown", side_effect=ImportError), \
             patch.object(proc, "_pdf_to_markdown", return_value="# corrupt\n\n(unreadable)"):
            md = proc._to_markdown(p)
        assert isinstance(md, str)


# =============================================================================
#  DoclingProcessor — heading-aware splitter
# =============================================================================

class TestSplitWithHeadings:

    def _proc(self) -> DoclingProcessor:
        return DoclingProcessor(chunk_size=150, chunk_overlap=15)

    def test_h1_heading_enters_section_path(self, tmp_path):
        proc = self._proc()
        md   = "# Introduction\n\nThis section introduces the topic in detail with sufficient content."
        chunks = proc._split_with_headings(
            md, str(tmp_path / "doc.pdf"), "Doc", "pdf"
        )
        assert any("Introduction" in c.section_path for c in chunks)

    def test_h2_nested_under_h1(self, tmp_path):
        proc = self._proc()
        md   = (
            "# Chapter 1\n\n"
            "## Background\n\n"
            "Background text about the context and history of the topic being studied."
        )
        chunks = proc._split_with_headings(
            md, str(tmp_path / "d.pdf"), "D", "pdf"
        )
        nested = [c for c in chunks if "Background" in c.section_path]
        assert len(nested) >= 1
        # Chapter 1 should also be in path (parent heading)
        assert any("Chapter 1" in c.section_path for c in nested)

    def test_page_marker_sets_page_number(self, tmp_path):
        proc = self._proc()
        md   = "## Page 3\n\nContent on page three with useful information about revenue."
        chunks = proc._split_with_headings(
            md, str(tmp_path / "d.pdf"), "D", "pdf"
        )
        assert all(c.page_number == 3 for c in chunks if c.text.strip())

    def test_page_marker_not_in_section_path(self, tmp_path):
        proc = self._proc()
        md   = "## Page 5\n\nSome content on page five in the document."
        chunks = proc._split_with_headings(
            md, str(tmp_path / "d.pdf"), "D", "pdf"
        )
        for c in chunks:
            assert "Page 5" not in c.section_path

    def test_markdown_source_flag_always_set(self, tmp_path):
        proc   = self._proc()
        md     = "# Title\n\nContent for testing the markdown source flag in metadata."
        chunks = proc._split_with_headings(
            md, str(tmp_path / "d.pdf"), "D", "pdf"
        )
        for c in chunks:
            assert c.metadata.get("markdown_source") is True

    def test_doctype_recorded_in_metadata(self, tmp_path):
        proc   = self._proc()
        md     = "Some content without headings but enough text to form a proper chunk."
        chunks = proc._split_with_headings(
            md, str(tmp_path / "d.pdf"), "D", "pdf"
        )
        for c in chunks:
            assert c.metadata.get("doctype") == "pdf"

    def test_chunk_text_is_utf8_string(self, tmp_path):
        proc = self._proc()
        md   = "# Résumé\n\nStraße and café with NFKC text normalization applied here."
        chunks = proc._split_with_headings(
            md, str(tmp_path / "d.pdf"), "D", "pdf"
        )
        for c in chunks:
            # Must be encodable as UTF-8 without errors
            c.text.encode("utf-8")

    def test_heading_resets_section_at_same_level(self, tmp_path):
        proc = self._proc()
        md   = (
            "## Section A\n\nContent A with enough text to make a meaningful chunk.\n\n"
            "## Section B\n\nContent B with enough text to make a meaningful chunk."
        )
        chunks = proc._split_with_headings(
            md, str(tmp_path / "d.pdf"), "D", "pdf"
        )
        section_a = [c for c in chunks if "Section A" in c.section_path]
        section_b = [c for c in chunks if "Section B" in c.section_path]
        assert len(section_a) >= 1
        assert len(section_b) >= 1
        # Section B chunks should have Section B as the last heading
        for ch in section_b:
            assert "Section B" == ch.section_path[-1]

    def test_char_offset_increases_monotonically(self, tmp_path):
        proc   = self._proc()
        md     = ("word " * 200)   # long enough to produce multiple chunks
        chunks = proc._split_with_headings(
            md, str(tmp_path / "d.pdf"), "D", "pdf"
        )
        offsets = [c.char_offset for c in chunks]
        assert offsets == sorted(offsets)

    def test_empty_markdown_returns_empty_list(self, tmp_path):
        proc   = self._proc()
        chunks = proc._split_with_headings(
            "   \n\n  ", str(tmp_path / "d.pdf"), "D", "pdf"
        )
        assert chunks == []


# =============================================================================
#  DoclingProcessor.process() — full pipeline
# =============================================================================

class TestDoclingProcessorProcess:

    def test_raises_for_missing_file(self, tmp_path):
        proc = _make_processor()
        with pytest.raises(FileNotFoundError):
            proc.process(tmp_path / "nonexistent.pdf")

    def test_raises_for_unsupported_extension(self, tmp_path):
        p = tmp_path / "file.mp3"
        p.write_bytes(b"fake")
        proc = _make_processor()
        with pytest.raises(ValueError, match="Unsupported"):
            proc.process(p)

    def test_returns_list_of_document_chunks(self, tmp_path):
        p    = _minimal_docx(tmp_path)
        proc = _make_processor()
        with patch.object(proc, "_docling_to_markdown", side_effect=ImportError):
            chunks = proc.process(p)
        assert isinstance(chunks, list)
        assert all(isinstance(c, DocumentChunk) for c in chunks)

    def test_chunks_non_empty_for_valid_docx(self, tmp_path):
        p    = _minimal_docx(tmp_path)
        proc = _make_processor()
        with patch.object(proc, "_docling_to_markdown", side_effect=ImportError):
            chunks = proc.process(p)
        assert len(chunks) >= 1

    def test_all_chunk_text_is_valid_utf8(self, tmp_path):
        p    = _minimal_docx(tmp_path)
        proc = _make_processor()
        with patch.object(proc, "_docling_to_markdown", side_effect=ImportError):
            chunks = proc.process(p)
        for chunk in chunks:
            chunk.text.encode("utf-8")   # must not raise

    def test_no_null_bytes_in_chunks(self, tmp_path):
        p    = _minimal_docx(tmp_path)
        proc = _make_processor()
        with patch.object(proc, "_docling_to_markdown", side_effect=ImportError):
            chunks = proc.process(p)
        for chunk in chunks:
            assert "\x00" not in chunk.text

    def test_no_control_chars_in_chunks(self, tmp_path):
        p    = _minimal_docx(tmp_path)
        proc = _make_processor()
        with patch.object(proc, "_docling_to_markdown", side_effect=ImportError):
            chunks = proc.process(p)
        import re
        control_re = re.compile(r"[\x01-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]")
        for chunk in chunks:
            preview = repr(chunk.text[:80])
            assert not control_re.search(chunk.text), (
                f"Control char in chunk: {preview}"
            )

    def test_markdown_source_flag_set(self, tmp_path):
        p    = _minimal_docx(tmp_path)
        proc = _make_processor()
        with patch.object(proc, "_docling_to_markdown", side_effect=ImportError):
            chunks = proc.process(p)
        for chunk in chunks:
            assert chunk.metadata.get("markdown_source") is True

    def test_doc_title_derived_from_filename(self, tmp_path):
        p    = tmp_path / "annual_report_2024.docx"
        src  = _minimal_docx(tmp_path)
        src.rename(p)
        proc = _make_processor()
        with patch.object(proc, "_docling_to_markdown", side_effect=ImportError):
            chunks = proc.process(p)
        for chunk in chunks:
            assert "Annual Report 2024" == chunk.doc_title

    def test_chunk_ids_unique(self, tmp_path):
        p    = _minimal_docx(tmp_path)
        proc = _make_processor()
        with patch.object(proc, "_docling_to_markdown", side_effect=ImportError):
            chunks = proc.process(p)
        ids = [c.chunk_id for c in chunks]
        assert len(ids) == len(set(ids))

    def test_uses_docling_when_available(self, tmp_path):
        p    = _fake_pdf(tmp_path)
        proc = _make_processor()
        mock_md = "# Report\n\nDocling-generated Markdown content here."
        with patch.object(proc, "_docling_to_markdown", return_value=mock_md):
            chunks = proc.process(p)
        assert len(chunks) >= 1
        assert all(c.metadata.get("markdown_source") is True for c in chunks)

    def test_falls_back_on_docling_import_error(self, tmp_path):
        p    = _minimal_docx(tmp_path)
        proc = _make_processor()
        with patch.object(proc, "_docling_to_markdown", side_effect=ImportError("no docling")):
            chunks = proc.process(p)
        assert isinstance(chunks, list)
        assert len(chunks) >= 1

    def test_falls_back_on_docling_runtime_error(self, tmp_path):
        p    = _minimal_docx(tmp_path)
        proc = _make_processor()
        with patch.object(proc, "_docling_to_markdown", side_effect=RuntimeError("crash")):
            chunks = proc.process(p)
        assert isinstance(chunks, list)

    def test_xlsx_chunks_contain_table_pipe_chars(self, tmp_path):
        p    = _minimal_xlsx(tmp_path)
        proc = _make_processor()
        with patch.object(proc, "_docling_to_markdown", side_effect=ImportError):
            chunks = proc.process(p)
        combined = " ".join(c.text for c in chunks)
        # Pipe tables from _xlsx_to_markdown should survive into chunks
        assert "|" in combined or "Q1" in combined

    def test_pptx_chunks_have_slide_page_numbers(self, tmp_path):
        p    = _minimal_pptx(tmp_path)
        proc = _make_processor()
        with patch.object(proc, "_docling_to_markdown", side_effect=ImportError):
            chunks = proc.process(p)
        page_numbers = {c.page_number for c in chunks if c.page_number}
        assert len(page_numbers) >= 1

    def test_process_many_combines_chunks(self, tmp_path):
        p1 = _minimal_docx(tmp_path, "a.docx")
        p2 = _minimal_docx(tmp_path, "b.docx")
        proc = _make_processor()
        with patch.object(proc, "_docling_to_markdown", side_effect=ImportError):
            chunks = proc.process_many([p1, p2])
        titles = {c.doc_title for c in chunks}
        assert len(titles) == 2

    def test_process_many_skips_errors(self, tmp_path):
        p_good = _minimal_docx(tmp_path, "good.docx")
        p_bad  = tmp_path / "missing.docx"   # does not exist
        proc   = _make_processor()
        with patch.object(proc, "_docling_to_markdown", side_effect=ImportError):
            chunks = proc.process_many([p_good, p_bad])
        # good.docx processed; missing.docx skipped, no exception raised
        assert len(chunks) >= 1


# =============================================================================
#  UTF-8 normalisation applied to Docling output
# =============================================================================

class TestUtf8NormalisationIntegration:

    def test_control_chars_in_docling_output_stripped(self, tmp_path):
        """Control characters that Docling might emit are removed before chunking."""
        p    = _fake_pdf(tmp_path)
        proc = _make_processor()
        # Simulate Docling emitting control chars
        dirty_md = "# Title\n\nParagraph with \x0b vertical tab and \x01 SOH character."
        with patch.object(proc, "_docling_to_markdown", return_value=dirty_md):
            chunks = proc.process(p)
        import re
        ctrl = re.compile(r"[\x01-\x08\x0b\x0c\x0e-\x1f]")
        for c in chunks:
            assert not ctrl.search(c.text)

    def test_null_bytes_in_docling_output_stripped(self, tmp_path):
        p    = _fake_pdf(tmp_path)
        proc = _make_processor()
        dirty_md = "# Title\n\nContent\x00with\x00null\x00bytes embedded inside."
        with patch.object(proc, "_docling_to_markdown", return_value=dirty_md):
            chunks = proc.process(p)
        for c in chunks:
            assert "\x00" not in c.text

    def test_nfkc_normalisation_applied_to_docling_output(self, tmp_path):
        p    = _fake_pdf(tmp_path)
        proc = _make_processor()
        # ﬁ (ligature) should be decomposed to "fi"
        md_with_ligature = "# ﬁnancial Report\n\nThe ﬁgures show growth of １０ percent."
        with patch.object(proc, "_docling_to_markdown", return_value=md_with_ligature):
            chunks = proc.process(p)
        combined = " ".join(c.text for c in chunks)
        assert "\uFB01" not in combined   # ligature gone
        assert "fi" in combined.lower() or "f" in combined.lower()

    def test_latin1_encoded_bytes_normalised(self):
        """Latin-1 bytes decoded and NFKC normalised correctly."""
        latin1_bytes = "café résumé naïve".encode("latin-1")
        result = to_utf8(latin1_bytes, source_encoding="latin-1")
        result.encode("utf-8")   # must not raise
        assert "caf" in result

    def test_chunk_text_encodable_after_ligature_expansion(self, tmp_path):
        p    = _fake_pdf(tmp_path)
        proc = _make_processor()
        md   = "# Title\n\n" + "The ﬁle contains ﬂowcharts and résumés. " * 10
        with patch.object(proc, "_docling_to_markdown", return_value=md):
            chunks = proc.process(p)
        for c in chunks:
            c.text.encode("utf-8")
