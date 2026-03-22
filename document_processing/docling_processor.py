"""
document_processing/docling_processor.py
──────────────────────────────────────────
Converts documents (PDF, DOCX, XLSX, PPTX) into structured text chunks
ready for vector-store ingestion.

Processing pipeline (per file)
───────────────────────────────
  Stage 1 — Parse → Markdown
    Docling converts the source file to a clean Markdown string that
    preserves headings, lists, tables (GFM pipe-table format), and code
    blocks.  This gives the embedder a semantically rich, structured view.

    When Docling is unavailable, lightweight pure-Python renderers produce
    an equivalent Markdown string:
        DOCX  →  python-docx  (heading styles → # / ## / ###)
        XLSX  →  openpyxl     (sheets → ## heading + pipe-tables)
        PPTX  →  python-pptx  (slides → ## Slide N + paragraphs)
        PDF   →  pypdf         (pages  → ## Page N + plain text)

  Stage 2 — UTF-8 normalisation  (``to_utf8``)
    • Decodes bytes to str (auto-detects encoding, errors→replacement char).
    • Round-trips through UTF-8 to flush lone surrogate code points.
    • Applies NFKC normalisation (ligatures, fullwidth digits, etc.).
    • Strips null bytes and C0/C1 control characters that break tokenisers.

  Stage 3 — Chunk with heading context
    The normalised Markdown is parsed line by line to track the active
    heading breadcrumb (# / ## / ###).  Content between headings is split
    into overlapping windows using chunk_size / chunk_overlap.  Every chunk
    carries its section_path so retrieval can surface provenance.
"""
from __future__ import annotations

import hashlib
import logging
import re
import unicodedata
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable, Optional

from config import settings

logger = logging.getLogger(__name__)

SUPPORTED_SUFFIXES = {".pdf", ".docx", ".xlsx", ".xls", ".pptx", ".ppt"}

_HEADING_RE  = re.compile(r"^(#{1,6})\s+(.*)")
_PAGE_RE     = re.compile(r"^#{1,2}\s+(?:Page|Slide)\s+(\d+)", re.IGNORECASE)
_CONTROL_RE  = re.compile(r"[\x01-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]")


# =============================================================================
#  UTF-8 normalisation
# =============================================================================

def to_utf8(
        text: "str | bytes",
        source_encoding: str = "utf-8",
        normalisation_form: str = "NFKC",
) -> str:
    """
    Return a clean, normalised UTF-8 string.

    Steps
    -----
    1. Decode ``bytes`` → ``str`` using ``source_encoding``
       (``errors="replace"`` so no byte is silently dropped).
    2. Round-trip through UTF-8 to flush lone surrogate code points.
    3. Apply Unicode normalisation (default NFKC):
       ﬁ→fi, １→1, ™→TM, etc.
    4. Strip null bytes and C0/C1 control characters
       (keeps \\t, \\n, \\r).

    Parameters
    ----------
    text              : Input str or bytes.
    source_encoding   : Codec used when *text* is bytes.
    normalisation_form: Unicode form — NFC | NFKC | NFD | NFKD.

    Returns
    -------
    str  Clean UTF-8 string safe for embedding models.
    """
    if isinstance(text, bytes):
        text = text.decode(source_encoding, errors="replace")

    # Flush lone surrogates that Python str can hold but UTF-8 cannot encode
    text = text.encode("utf-8", errors="replace").decode("utf-8", errors="replace")

    # Canonical / compatibility decomposition + recomposition
    text = unicodedata.normalize(normalisation_form, text)

    # Replace control chars (except \t \n \r) with a space
    text = _CONTROL_RE.sub(" ", text)

    # Remove null bytes
    text = text.replace("\x00", "")

    return text


# =============================================================================
#  DocumentChunk
# =============================================================================

@dataclass
class DocumentChunk:
    """
    A single Markdown text chunk extracted from a document.

    All ``text`` content is guaranteed to be valid UTF-8 Markdown.

    Attributes
    ----------
    chunk_id      : SHA-256 derived ID (source path + page + offset).
    text          : UTF-8 Markdown content for this chunk.
    doc_path      : Absolute path to the source document.
    doc_title     : Human-readable document name.
    page_number   : 1-based page / slide / sheet number (None if unknown).
    section_path  : Heading breadcrumb, e.g. ["Chapter 1", "Introduction"].
    char_offset   : Character offset within the document (ordering key).
    metadata      : Arbitrary extra metadata.
                    ``markdown_source=True`` is always set so consumers know
                    the text is Markdown-formatted.
    """

    chunk_id:     str
    text:         str
    doc_path:     str
    doc_title:    str
    page_number:  Optional[int] = None
    section_path: list[str]     = field(default_factory=list)
    char_offset:  int           = 0
    metadata:     dict          = field(default_factory=dict)

    @property
    def reference(self) -> str:
        """'Q3_Report › Page 4 › Revenue Analysis'."""
        parts = [self.doc_title]
        if self.page_number:
            parts.append(f"Page {self.page_number}")
        if self.section_path:
            parts.append(" › ".join(self.section_path))
        return " › ".join(parts)

    def to_langchain_doc(self):
        """Convert to a LangChain Document for vector-store ingestion."""
        from langchain_core.documents import Document
        return Document(
            page_content=self.text,
            metadata={
                "chunk_id":     self.chunk_id,
                "doc_path":     self.doc_path,
                "doc_title":    self.doc_title,
                "page_number":  self.page_number,
                "section_path": " › ".join(self.section_path),
                "reference":    self.reference,
                **self.metadata,
            },
        )


# =============================================================================
#  DoclingProcessor
# =============================================================================

class DoclingProcessor:
    """
    Convert documents to UTF-8 Markdown chunks.

    Per-file pipeline
    -----------------
    1. ``_to_markdown(path)``      Docling or fallback renderer → raw Markdown
    2. ``to_utf8(markdown)``       Normalise encoding + strip control chars
    3. ``_split_with_headings()``  Track heading context; emit DocumentChunks
    """

    def __init__(
            self,
            chunk_size:    int = settings.chunk_size,
            chunk_overlap: int = settings.chunk_overlap,
    ) -> None:
        self.chunk_size    = chunk_size
        self.chunk_overlap = chunk_overlap

    # -------------------------------------------------------------------------
    #  Helpers
    # -------------------------------------------------------------------------

    @staticmethod
    def _make_chunk_id(doc_path: str, page: Optional[int], offset: int) -> str:
        raw = f"{doc_path}:{page}:{offset}"
        return hashlib.sha256(raw.encode()).hexdigest()[:16]

    def _split_text(self, text: str) -> list[str]:
        """Overlapping sliding-window splitter."""
        chunks: list[str] = []
        step = max(1, self.chunk_size - self.chunk_overlap)
        for start in range(0, len(text), step):
            chunk = text[start : start + self.chunk_size]
            if chunk.strip():
                chunks.append(chunk)
        return chunks

    # -------------------------------------------------------------------------
    #  Stage 1 — Parse → Markdown
    # -------------------------------------------------------------------------

    def _to_markdown(self, path: Path) -> str:
        """
        Convert *path* to a Markdown string.

        Tries Docling first (best quality — OCR, tables, heading detection).
        Falls back to pure-Python renderers when Docling is absent or fails.
        """
        try:
            return self._docling_to_markdown(path)
        except ImportError:
            logger.warning(
                "Docling not available — using fallback Markdown renderer for '%s'.",
                path.name,
            )
        except Exception as exc:
            logger.warning(
                "Docling failed for '%s' (%s) — using fallback renderer.",
                path.name, exc,
            )
        return self._fallback_to_markdown(path)

    def _docling_to_markdown(self, path: Path) -> str:
        """
        Use Docling's DocumentConverter to produce Markdown.

        Docling's ``export_to_markdown()`` preserves:
          • Headings (# / ## / ###)
          • Paragraphs and lists
          • Tables (GFM pipe-table format)
          • Code blocks (fenced with ```)
          • Captions and footnotes
        """
        try:
            from docling.document_converter import DocumentConverter  # type: ignore
        except ImportError as exc:
            raise ImportError(
                "docling is not installed.  Run: pip install docling"
            ) from exc

        converter = DocumentConverter()
        result    = converter.convert(str(path))
        markdown  = result.document.export_to_markdown()
        logger.debug(
            "Docling→Markdown: '%s' → %d chars.", path.name, len(markdown)
        )
        return markdown

    def _fallback_to_markdown(self, path: Path) -> str:
        """Dispatch to the correct pure-Python Markdown renderer."""
        suffix = path.suffix.lower()
        if suffix == ".docx":
            return self._docx_to_markdown(path)
        if suffix in (".xlsx", ".xls"):
            return self._xlsx_to_markdown(path)
        if suffix in (".pptx", ".ppt"):
            return self._pptx_to_markdown(path)
        # .pdf or unknown
        return self._pdf_to_markdown(path)

    def _docx_to_markdown(self, path: Path) -> str:
        import docx  # type: ignore
        document = docx.Document(str(path))
        blocks: list[str] = []
        for para in document.paragraphs:
            text  = para.text.strip()
            if not text:
                continue
            style = (para.style.name or "") if para.style else ""
            if style.startswith("Heading"):
                try:
                    level = int(style.split()[-1])
                except ValueError:
                    level = 1
                blocks.append(f"{'#' * min(level, 6)} {text}")
            else:
                blocks.append(text)
        return "\n\n".join(blocks)

    def _xlsx_to_markdown(self, path: Path) -> str:
        import openpyxl  # type: ignore
        wb     = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        blocks: list[str] = []
        for sheet_name in wb.sheetnames:
            ws   = wb[sheet_name]
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue
            blocks.append(f"## {sheet_name}")
            header = [str(c) if c is not None else "" for c in rows[0]]
            sep    = ["---"] * len(header)
            table  = ["| " + " | ".join(header) + " |",
                      "| " + " | ".join(sep)    + " |"]
            for row in rows[1:]:
                cells = [str(c) if c is not None else "" for c in row]
                while len(cells) < len(header):
                    cells.append("")
                table.append("| " + " | ".join(cells[: len(header)]) + " |")
            blocks.append("\n".join(table))
        wb.close()
        return "\n\n".join(blocks)

    def _pptx_to_markdown(self, path: Path) -> str:
        from pptx import Presentation  # type: ignore
        prs    = Presentation(str(path))
        blocks: list[str] = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                blocks.append(f"## Slide {slide_num}")
                blocks.append("\n\n".join(texts))
        return "\n\n".join(blocks)

    def _pdf_to_markdown(self, path: Path) -> str:
        """
        Last-resort text extraction from PDF using pypdf.
        Produces ``## Page N`` headings with plain text paragraphs.
        """
        try:
            import pypdf  # type: ignore
            reader = pypdf.PdfReader(str(path))
            pages: list[str] = []
            for i, page in enumerate(reader.pages, start=1):
                text = (page.extract_text() or "").strip()
                if text:
                    pages.append(f"## Page {i}\n\n{text}")
            if pages:
                return "\n\n".join(pages)
        except ImportError:
            logger.debug("pypdf not installed; skipping PDF text extraction.")
        except Exception as exc:
            logger.debug("pypdf failed for '%s': %s", path.name, exc)

        # Absolute last resort — raw byte decode
        try:
            raw  = path.read_bytes()
            text = raw.decode("utf-8", errors="replace")
            text = re.sub(r"[^\x09\x0A\x0D\x20-\x7E\u0080-\uFFFF]+", " ", text)
            return text[:50_000]
        except Exception:
            pass
        return f"# {path.stem}\n\n(Could not extract text from this PDF.)"

    # -------------------------------------------------------------------------
    #  Stage 3 — Heading-aware splitter
    # -------------------------------------------------------------------------

    def _split_with_headings(
            self,
            markdown:    str,
            doc_path_str: str,
            doc_title:   str,
            doctype:     str,
    ) -> list[DocumentChunk]:
        """
        Parse Markdown line-by-line, tracking the active heading stack,
        then flush buffered content into overlapping chunks annotated with
        the current section_path and page_number.

        ``## Page N`` / ``## Slide N`` markers (emitted by fallback renderers)
        update ``current_page`` but are NOT added to the section breadcrumb.
        All other headings update the breadcrumb stack normally.
        """
        lines:            list[str]          = markdown.splitlines()
        current_headings: list[str]          = []
        current_page:     Optional[int]      = None
        buffer:           list[str]          = []
        chunks:           list[DocumentChunk] = []
        global_offset:    int                = 0

        def _flush(buf: list[str], hdgs: list[str], page: Optional[int]) -> None:
            nonlocal global_offset
            joined = "\n".join(buf).strip()
            if not joined:
                return
            for sub in self._split_text(joined):
                sub = to_utf8(sub)          # re-normalise every chunk
                if not sub.strip():
                    continue
                cid = self._make_chunk_id(doc_path_str, page, global_offset)
                chunks.append(DocumentChunk(
                    chunk_id=cid,
                    text=sub,
                    doc_path=doc_path_str,
                    doc_title=doc_title,
                    page_number=page,
                    section_path=list(hdgs),
                    char_offset=global_offset,
                    metadata={"doctype": doctype, "markdown_source": True},
                ))
                global_offset += len(sub)

        for line in lines:
            hm = _HEADING_RE.match(line)
            if hm:
                _flush(buffer, current_headings, current_page)
                buffer = []

                level   = len(hm.group(1))
                heading = hm.group(2).strip()[:120]

                pm = _PAGE_RE.match(line)
                if pm:
                    # "## Page 3" / "## Slide 3" — page marker from fallback renderer
                    current_page = int(pm.group(1))
                else:
                    # Normal heading — update breadcrumb stack
                    current_headings = current_headings[: level - 1]
                    current_headings.append(heading)
            else:
                buffer.append(line)

        _flush(buffer, current_headings, current_page)
        return chunks

    # -------------------------------------------------------------------------
    #  Public API
    # -------------------------------------------------------------------------

    def process(self, path: "str | Path") -> list[DocumentChunk]:
        """
        Parse *path* and return UTF-8 Markdown chunks.

        Pipeline
        --------
        1. Convert to Markdown  (Docling ▸ fallback renderer).
        2. Normalise to UTF-8   (NFKC, strip control chars).
        3. Split into chunks    (heading-aware overlapping windows).

        Parameters
        ----------
        path : str | Path

        Returns
        -------
        list[DocumentChunk]  In document order; empty if no text found.

        Raises
        ------
        FileNotFoundError  – file does not exist.
        ValueError         – unsupported file type.
        """
        path = Path(path)
        if not path.exists():
            raise FileNotFoundError(f"Document not found: {path}")
        if path.suffix.lower() not in SUPPORTED_SUFFIXES:
            raise ValueError(
                f"Unsupported file type '{path.suffix}'. "
                f"Allowed: {', '.join(sorted(SUPPORTED_SUFFIXES))}"
            )

        doc_path_str = str(path.resolve())
        doc_title    = path.stem.replace("_", " ").replace("-", " ").title()
        doctype      = path.suffix.lstrip(".")

        # Stage 1 — Parse → Markdown
        raw_markdown = self._to_markdown(path)

        # Stage 2 — UTF-8 normalisation
        normalised = to_utf8(raw_markdown)

        if not normalised.strip():
            logger.warning("'%s' produced no extractable text.", path.name)
            return []

        # Stage 3 — Chunk with heading context
        chunks = self._split_with_headings(
            normalised, doc_path_str, doc_title, doctype
        )

        logger.info(
            "Processed '%s' → %d Markdown chunks (%.1f KB).",
            path.name, len(chunks), len(normalised) / 1024,
                                    )
        return chunks

    def process_many(self, paths: "Iterable[str | Path]") -> list[DocumentChunk]:
        """Process multiple documents; return all chunks combined."""
        all_chunks: list[DocumentChunk] = []
        for p in paths:
            try:
                all_chunks.extend(self.process(p))
            except Exception as exc:
                logger.error("Failed to process '%s': %s", p, exc)
        return all_chunks