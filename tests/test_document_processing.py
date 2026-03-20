"""
tests/test_document_processing.py
───────────────────────────────────
Unit tests for:
  • DoclingProcessor (chunk generation, fallback parsers, metadata)
  • VectorStore      (ingest, search, dedup, document-level ops)
  • DocumentManager  (façade behaviour)

Uses temporary directories and in-memory fixtures to avoid hitting
the real vector store or Docling inference.
"""
from __future__ import annotations

import hashlib
import tempfile
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest

from document_processing.docling_processor import (
    DocumentChunk,
    DoclingProcessor,
    SUPPORTED_SUFFIXES,
)


# ─────────────────────────────────────────────────────────────────────────────
#  Fixtures
# ─────────────────────────────────────────────────────────────────────────────

@pytest.fixture()
def tmp_dir(tmp_path: Path) -> Path:
    return tmp_path


@pytest.fixture()
def plain_txt(tmp_dir: Path) -> Path:
    """A plain-text file used as a stand-in for unsupported format testing."""
    p = tmp_dir / "sample.txt"
    p.write_text("Hello world. This is a test document.\nSecond paragraph here.")
    return p


@pytest.fixture()
def sample_docx(tmp_dir: Path) -> Path:
    """Create a minimal DOCX file for fallback parser testing."""
    try:
        import docx

        doc = docx.Document()
        doc.add_heading("Introduction", level=1)
        doc.add_paragraph("This is the introduction paragraph.")
        doc.add_heading("Details", level=2)
        doc.add_paragraph("Some detailed content for section two.")
        path = tmp_dir / "sample.docx"
        doc.save(str(path))
        return path
    except ImportError:
        pytest.skip("python-docx not installed")


@pytest.fixture()
def sample_xlsx(tmp_dir: Path) -> Path:
    """Create a minimal XLSX file for fallback parser testing."""
    try:
        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Sales"
        ws.append(["Month", "Revenue", "Units"])
        ws.append(["January", 50000, 120])
        ws.append(["February", 62000, 145])
        path = tmp_dir / "sample.xlsx"
        wb.save(str(path))
        return path
    except ImportError:
        pytest.skip("openpyxl not installed")


@pytest.fixture()
def sample_pptx(tmp_dir: Path) -> Path:
    """Create a minimal PPTX file for fallback parser testing."""
    try:
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Test Slide"
        slide.placeholders[1].text = "Slide content goes here."
        path = tmp_dir / "sample.pptx"
        prs.save(str(path))
        return path
    except ImportError:
        pytest.skip("python-pptx not installed")


@pytest.fixture()
def processor() -> DoclingProcessor:
    return DoclingProcessor(chunk_size=100, chunk_overlap=20)


# ─────────────────────────────────────────────────────────────────────────────
#  DocumentChunk tests
# ─────────────────────────────────────────────────────────────────────────────

class TestDocumentChunk:
    def test_reference_full(self):
        chunk = DocumentChunk(
            chunk_id="abc123",
            text="Some text.",
            doc_path="/docs/report.pdf",
            doc_title="Report",
            page_number=3,
            section_path=["Chapter 1", "Revenue"],
        )
        ref = chunk.reference
        assert "Report" in ref
        assert "Page 3" in ref
        assert "Chapter 1" in ref
        assert "Revenue" in ref

    def test_reference_minimal(self):
        chunk = DocumentChunk(
            chunk_id="abc123",
            text="Hello",
            doc_path="/docs/report.pdf",
            doc_title="Report",
        )
        assert chunk.reference == "Report"

    def test_to_langchain_doc(self):
        chunk = DocumentChunk(
            chunk_id="xyz",
            text="Test content.",
            doc_path="/tmp/doc.pdf",
            doc_title="Doc",
            page_number=1,
            section_path=["Intro"],
        )
        lc_doc = chunk.to_langchain_doc()
        assert lc_doc.page_content == "Test content."
        assert lc_doc.metadata["chunk_id"] == "xyz"
        assert lc_doc.metadata["page_number"] == 1
        assert "Intro" in lc_doc.metadata["section_path"]


# ─────────────────────────────────────────────────────────────────────────────
#  DoclingProcessor tests
# ─────────────────────────────────────────────────────────────────────────────

class TestDoclingProcessor:

    def test_supported_suffixes(self):
        assert ".pdf" in SUPPORTED_SUFFIXES
        assert ".docx" in SUPPORTED_SUFFIXES
        assert ".xlsx" in SUPPORTED_SUFFIXES
        assert ".pptx" in SUPPORTED_SUFFIXES

    def test_unsupported_file_raises(self, processor: DoclingProcessor, tmp_dir: Path):
        p = tmp_dir / "bad.csv"
        p.write_text("a,b,c\n1,2,3")
        with pytest.raises(ValueError, match="Unsupported"):
            processor.process(p)

    def test_missing_file_raises(self, processor: DoclingProcessor):
        with pytest.raises(FileNotFoundError):
            processor.process("/nonexistent/path/file.pdf")

    def test_text_splitter(self, processor: DoclingProcessor):
        text = "A" * 250
        chunks = processor._split_text(text)
        assert len(chunks) >= 2
        assert all(len(c) <= processor.chunk_size for c in chunks)

    def test_chunk_id_stable(self, processor: DoclingProcessor):
        cid1 = processor._make_chunk_id("/path/doc.pdf", 1, 0)
        cid2 = processor._make_chunk_id("/path/doc.pdf", 1, 0)
        assert cid1 == cid2
        assert len(cid1) == 16

    def test_fallback_docx(self, processor: DoclingProcessor, sample_docx: Path):
        with patch.object(processor, "_parse_with_docling", side_effect=ImportError("no docling")):
            chunks = processor.process(sample_docx)
        assert len(chunks) > 0
        combined = " ".join(c.text for c in chunks)
        assert "introduction" in combined.lower() or "paragraph" in combined.lower()

    def test_fallback_xlsx(self, processor: DoclingProcessor, sample_xlsx: Path):
        with patch.object(processor, "_parse_with_docling", side_effect=ImportError("no docling")):
            chunks = processor.process(sample_xlsx)
        assert len(chunks) > 0
        assert any("Sales" in c.section_path for c in chunks if c.section_path)

    def test_fallback_pptx(self, processor: DoclingProcessor, sample_pptx: Path):
        with patch.object(processor, "_parse_with_docling", side_effect=ImportError("no docling")):
            chunks = processor.process(sample_pptx)
        assert len(chunks) > 0
        assert all(c.page_number is not None for c in chunks)


# ─────────────────────────────────────────────────────────────────────────────
#  VectorStore tests (mocked embeddings + chroma)
# ─────────────────────────────────────────────────────────────────────────────

class TestVectorStore:
    """Tests that mock out ChromaDB and the embedding model."""

    def _make_chunks(self, n: int = 3, doc_title: str = "Test Doc") -> list[DocumentChunk]:
        return [
            DocumentChunk(
                chunk_id=f"chunk_{i:04d}",
                text=f"Chunk {i} content about {doc_title}.",
                doc_path=f"/docs/{doc_title}.pdf",
                doc_title=doc_title,
                page_number=i + 1,
                section_path=[f"Section {i}"],
            )
            for i in range(n)
        ]

    def test_ingest_new_chunks(self, tmp_dir: Path):
        from document_processing.vector_store import VectorStore

        mock_embeddings = MagicMock()
        mock_embeddings.embed_documents.return_value = [[0.1, 0.2, 0.3]] * 3
        mock_embeddings.embed_query.return_value = [0.1, 0.2, 0.3]

        mock_collection = MagicMock()
        mock_collection.count.return_value = 0
        mock_collection.get.return_value = {"ids": []}

        with (
            patch("document_processing.vector_store.get_embeddings", return_value=mock_embeddings),
            patch("chromadb.PersistentClient") as mock_client_cls,
        ):
            mock_client = MagicMock()
            mock_client.get_or_create_collection.return_value = mock_collection
            mock_client_cls.return_value = mock_client

            store = VectorStore(persist_directory=tmp_dir, collection_name="test")
            store._ensure_ready()
            chunks = self._make_chunks(3)
            added = store.ingest(chunks)

        assert added == 3
        mock_collection.add.assert_called_once()

    def test_ingest_deduplication(self, tmp_dir: Path):
        from document_processing.vector_store import VectorStore

        chunks = self._make_chunks(3)
        existing_ids = [c.chunk_id for c in chunks]

        mock_embeddings = MagicMock()
        mock_collection = MagicMock()
        mock_collection.count.return_value = 3
        mock_collection.get.return_value = {"ids": existing_ids}

        with (
            patch("document_processing.vector_store.get_embeddings", return_value=mock_embeddings),
            patch("chromadb.PersistentClient") as mock_client_cls,
        ):
            mock_client = MagicMock()
            mock_client.get_or_create_collection.return_value = mock_collection
            mock_client_cls.return_value = mock_client

            store = VectorStore(persist_directory=tmp_dir)
            store._ensure_ready()
            added = store.ingest(chunks)

        assert added == 0
        mock_collection.add.assert_not_called()


# ─────────────────────────────────────────────────────────────────────────────
#  DocumentManager integration tests (all dependencies mocked)
# ─────────────────────────────────────────────────────────────────────────────

class TestDocumentManager:

    def test_ingest_calls_processor_and_store(self, tmp_dir: Path, sample_docx: Path):
        from document_processing.document_manager import DocumentManager

        mock_chunks = [
            DocumentChunk("id1", "text1", str(sample_docx), "Sample", page_number=1)
        ]
        mock_store = MagicMock()
        mock_store.ingest.return_value = 1
        mock_store.count = 1

        with (
            patch.object(DoclingProcessor, "process", return_value=mock_chunks),
            patch("document_processing.document_manager.VectorStore", return_value=mock_store),
        ):
            dm = DocumentManager()
            count = dm.ingest(sample_docx)

        assert count == 1
        mock_store.ingest.assert_called_once_with(mock_chunks)

    def test_format_search_results_empty(self):
        from document_processing.document_manager import DocumentManager
        dm = DocumentManager()
        result = dm.format_search_results([])
        assert "No relevant" in result
